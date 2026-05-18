"""
rag_service.py — Lazy-loading RAG service for cloud deployment.

Models are loaded on first use (not at import time) to stay within
the 512MB RAM limit of Render's free tier.
"""

import zipfile
import xml.etree.ElementTree as ET
import re

def extract_ppt_text(file_path):
    """
    Memory-safe PPTX extraction.
    PPTX files are ZIP archives. We extract text directly from the XML
    to prevent python-pptx from loading large embedded images into RAM and causing OOM.
    """
    slides = []
    
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # Find all slide XML files
            slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            
            # Sort by slide number: slide1.xml, slide2.xml, etc.
            def get_slide_num(filename):
                match = re.search(r'slide(\d+)\.xml', filename)
                return int(match.group(1)) if match else 0
                
            slide_files.sort(key=get_slide_num)
            
            for i, slide_file in enumerate(slide_files):
                xml_content = z.read(slide_file)
                root = ET.fromstring(xml_content)
                
                # In DrawingML, text is inside <a:t> tags
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                
                texts = []
                for node in root.findall('.//a:t', ns):
                    if node.text:
                        texts.append(node.text)
                        
                slides.append({
                    "slide_number": i + 1,
                    "content": " ".join(texts).strip()
                })
                
        return slides
    except Exception as e:
        print(f"[PPTX Extraction Warning] Fast extraction failed: {e}. Falling back to python-pptx.")
        from pptx import Presentation
        presentation = Presentation(file_path)
        slides = []
        for i, slide in enumerate(presentation.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
            slides.append({
                "slide_number": i + 1,
                "content": text.strip()
            })
        return slides


def assess_ppt_quality(slides):
    """
    Inspect the extracted slides and return a quality report dict.
    Determines whether the PPT has meaningful content or is blank/poor.
    """
    total_slides = len(slides)
    blank_slides = [s for s in slides if len(s["content"].strip()) < 5]
    total_words = sum(len(s["content"].split()) for s in slides)
    non_blank = total_slides - len(blank_slides)

    has_content = total_words >= 20 and non_blank > 0

    if total_slides == 0:
        verdict = "NO_SLIDES"
        summary = "The uploaded file contains no slides at all. The PPT is completely empty."
    elif len(blank_slides) == total_slides:
        verdict = "ALL_BLANK"
        summary = (
            f"The uploaded PPT has {total_slides} slide(s) but ALL are blank with no text content. "
            "This is an unacceptable presentation file. There is nothing to evaluate from the PPT side."
        )
    elif total_words < 20:
        verdict = "EXTREMELY_POOR"
        summary = (
            f"The uploaded PPT has {total_slides} slide(s) with only {total_words} total words. "
            "This is an extremely poor and insufficient presentation with almost no content."
        )
    elif non_blank < 3:
        verdict = "POOR"
        summary = (
            f"The uploaded PPT has {total_slides} slide(s) but only {non_blank} contain actual text. "
            f"The total word count is {total_words}. This is a weak presentation."
        )
    else:
        verdict = "ACCEPTABLE"
        summary = (
            f"The PPT has {total_slides} slides with {total_words} total words across {non_blank} non-blank slides."
        )

    return {
        "verdict": verdict,
        "summary": summary,
        "total_slides": total_slides,
        "blank_slides": len(blank_slides),
        "total_words": total_words,
        "has_content": has_content
    }


def store_embeddings(session_id, slides):
    # DEPRECATED: We no longer store embeddings in ChromaDB to save memory.
    # Text is now injected directly into the ppt_quality dict.
    pass

def retrieve_relevant(query):
    # DEPRECATED: RAG is now handled by passing the full PPT text to Gemini.
    return ["[No relevant PPT content found]"]